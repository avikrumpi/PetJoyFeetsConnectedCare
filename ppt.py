import requests
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# --- CONFIGURATION ---
FILENAME = "PetLiv_AtoZ_Mart_AI_Presentation.pptx"

# Brand Colors
BLUE = RGBColor(37, 99, 235)    # #2563eb
CYAN = RGBColor(6, 182, 212)    # #06b6d4
ORANGE = RGBColor(249, 115, 22) # #f97316
DARK_GREY = RGBColor(30, 41, 59)# #1e293b
LIGHT_GREY = RGBColor(241, 245, 249) # #f1f5f9
WHITE = RGBColor(255, 255, 255)

# Image URLs (from our previous context)
IMG_TITLE = "http://googleusercontent.com/image_generation_content/15" 
IMG_SPLIT = "http://googleusercontent.com/image_generation_content/7"
IMG_VISION = "http://googleusercontent.com/image_collection/image_retrieval/4658569013286638380"
IMG_PROMISE = "http://googleusercontent.com/image_generation_content/1"
IMG_HYBRID = "http://googleusercontent.com/image_collection/image_retrieval/14454951342680244647"

# Presentation Dimensions (Fixed for the entire deck)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# --- HELPER FUNCTIONS ---

def add_header_strip(slide, prs_width):
    """Adds the vibrant top strip to a slide."""
    height = Cm(0.3)
    
    # Blue part (0 to 1/3)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs_width / 3, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = BLUE
    shape.line.fill.background()

    # Cyan part (1/3 to 2/3)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, prs_width / 3, 0, prs_width / 3, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = CYAN
    shape.line.fill.background()
    
    # Orange part (2/3 to end)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, (prs_width / 3) * 2, 0, prs_width / 3, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = ORANGE
    shape.line.fill.background()

def add_company_badge(slide):
    """Adds the company name badge."""
    textbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(4), Inches(0.5))
    tf = textbox.text_frame
    tf.margin_left = Cm(0.5)
    tf.margin_top = Cm(0.2)
    p = tf.add_paragraph()
    p.text = "Aavaaya HealthCare Pvt Limited"
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = BLUE
    
    # Add a rounded rectangle as background for the badge
    badge_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(0.25), Inches(4.2), Inches(0.6))
    badge_bg.fill.solid()
    badge_bg.fill.fore_color.rgb = RGBColor(239, 246, 255) # Light Blue BG
    badge_bg.line.color.rgb = BLUE
    badge_bg.line.width = Pt(1)
    slide.shapes[0].element.insert(1, badge_bg.element) # Move text box to the front
    
def fetch_image(url):
    """Downloads an image from a URL into a BytesIO object."""
    try:
        response = requests.get(url)
        response.raise_for_status()
        return BytesIO(response.content)
    except Exception as e:
        print(f"Could not download image {url}: {e}")
        return None

def set_text_style(paragraph, font_size, is_bold=False, color=DARK_GREY):
    paragraph.font.size = Pt(font_size)
    paragraph.font.bold = is_bold
    paragraph.font.color.rgb = color
    paragraph.font.name = 'Poppins'
    
def create_title_box(slide, text, size=36, color=BLUE, x=0.5, y=1, w=12, h=1):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    p = tb.text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = color
    p.font.name = 'Poppins'
    return tb

# --- MAIN SCRIPT ---

prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT

# --- SLIDE 1: Title Slide ---
slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank layout
add_header_strip(slide, prs.slide_width)
add_company_badge(slide)

# Title Text
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6), Inches(2))
tf = title_box.text_frame
tf.word_wrap = True
p = tf.add_paragraph()
p.text = "PetLiv AtoZ Mart AI"
set_text_style(p, 54, True, BLUE)

p = tf.add_paragraph()
p.text = "AI-Powered AtoZ Marketplace for Pets & Livestock – Digitize, Connect, Care, and Thrive Pan-India."
set_text_style(p, 24, False, DARK_GREY)

# Disclaimer Box
disclaimer_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.5), Inches(5.5), Inches(2.2))
disclaimer_box.fill.solid()
disclaimer_box.fill.fore_color.rgb = RGBColor(240, 249, 255) # Light Blue
disclaimer_box.line.color.rgb = CYAN
disclaimer_box.line.width = Pt(2)

tf = disclaimer_box.text_frame
tf.margin_left = Inches(0.2)
tf.margin_top = Inches(0.3)
p = tf.add_paragraph()
p.text = "Dedicated to Ethical Care:"
set_text_style(p, 14, True, BLUE)
p = tf.add_paragraph()
p.text = "We deliver happiness through loving connections, not sales.\nPetLiv strictly avoids selling/prescribing medicines, influencing clinical decisions, or enabling animal bidding. Our sole purpose is to enable trusted, ethical care for all Pet Parents and Livestock Caregivers."
set_text_style(p, 12, False, DARK_GREY)

# Image
img_stream = fetch_image(IMG_TITLE)
if img_stream:
    slide.shapes.add_picture(img_stream, Inches(7), Inches(1.2), Inches(6), Inches(5.5))


# --- SLIDE 2: Problem Statement ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header_strip(slide, prs.slide_width)
add_company_badge(slide)

# Left Text Area
tb = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(5.5), Inches(4))
tf = tb.text_frame
tf.word_wrap = True
p = tf.add_paragraph()
p.text = "The Challenge"
set_text_style(p, 36, True, ORANGE)

# Problem statement with Orange stripe effect
p_statement_text = '"Most Indian pet and livestock families face fragmented care, paperwork-heavy schemes, and poor digitization, blocking trusted vets, products, and benefits."'
p_statement_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(3.2), Inches(5.5), Inches(2.5))
p_statement_shape.fill.background()
p_statement_shape.line.fill.background()
p_statement_shape.left = Inches(0.5)
p_statement_shape.top = Inches(3.2)
p_statement_shape.width = Inches(5.5)
p_statement_shape.height = Inches(2.5)

tf_statement = p_statement_shape.text_frame
tf_statement.margin_left = Inches(0.2)
tf_statement.margin_top = Inches(0.2)
tf_statement.word_wrap = True
p = tf_statement.add_paragraph()
p.text = p_statement_text
set_text_style(p, 24, False, DARK_GREY)

# Add orange accent bar
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(3.2), Inches(0.1), Inches(2.5))
bar.fill.solid()
bar.fill.fore_color.rgb = ORANGE
bar.line.fill.background()


p = tf.add_paragraph()
p.text = "\nLet's Dive Deep.. ->"
set_text_style(p, 20, True, BLUE)
p.space_before = Pt(20)

# Image
img_stream = fetch_image(IMG_SPLIT)
if img_stream:
    slide.shapes.add_picture(img_stream, Inches(6.5), Inches(1.5), Inches(6.5), Inches(4.5))


# --- SLIDE 3: 9 Critical Ground Challenges (Grid) ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header_strip(slide, prs.slide_width)
add_company_badge(slide)

create_title_box(slide, "9 Critical Ground Challenges", x=0.5, y=0.8)

challenges = [
    ("Unified Digital Records", "78% lack records. Repeat diagnostics cost ₹8K/yr.", 0, BLUE),
    ("Offline & Invisible", "97% livestock undigitized.", 0, BLUE),
    ("Vet Access Barrier", "1:5,000 ratio. 45-min emergency waits.", 0, BLUE),
    ("One-Stop Marketplace", "15+ apps needed. 3-5 hours wasted weekly.", 1, CYAN),
    ("Counterfeit Supply", "30% fake meds cause 20% mortality.", 1, CYAN),
    ("Treatment Abandonment", "55% miss doses. ₹15K Cr losses.", 1, CYAN),
    ("Govt Scheme Access", "<20% tap funds due to paperwork.", 2, ORANGE),
    ("Rural Voice Exclusion", "70% rural women excluded by English apps.", 2, ORANGE),
    ("Language Gap", "90% content English-only.", 2, ORANGE)
]

row_y = [2.0, 3.8, 5.6]
col_x = [0.5, 4.7, 8.9]

for i, (head, body, r, color) in enumerate(challenges):
    c = i % 3
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(col_x[c]), Inches(row_y[r]), Inches(4), Inches(1.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_GREY
    shape.line.color.rgb = color
    shape.line.width = Pt(2)
    
    tf = shape.text_frame
    tf.margin_top = Cm(0.2)
    p = tf.add_paragraph()
    p.text = head
    set_text_style(p, 16, True, color)
    p = tf.add_paragraph()
    p.text = body
    set_text_style(p, 14, False, DARK_GREY)


# --- SLIDE 4: Why Now? (Text Overlay) ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header_strip(slide, prs.slide_width)
add_company_badge(slide)

# BG Shape (Dark)
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = DARK_GREY

# Text
tb = slide.shapes.add_textbox(Inches(2), Inches(2), Inches(9.333), Inches(4))
tf = tb.text_frame
tf.word_wrap = True
p = tf.add_paragraph()
p.text = "Why Now? The Urgency of Converging Deadlines"
set_text_style(p, 44, True, CYAN)
p.alignment = PP_ALIGN.CENTER

p = tf.add_paragraph()
p.text = "A perfect storm of government mandates (NDLM), explosive market growth (Pet Boom), and technological readiness (Agentic AI) creates a **must-act window** for digital disruption."
set_text_style(p, 24, False, WHITE)
p.alignment = PP_ALIGN.CENTER


# --- SLIDE 5: 6 Reasons ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header_strip(slide, prs.slide_width)
add_company_badge(slide)

create_title_box(slide, "6 Critical Reasons for Immediate Action", x=0.5, y=0.8)

reasons = [
    ("NDLM Mandate & Compliance", "100% ID by 2026. 574M animals risk missing DBT.", BLUE),
    ("Govt Funding Rush", "₹10K Cr unutilized funds. Window closing.", ORANGE),
    ("Explosive Pet Boom", "20% CAGR to ₹60K Cr by 2028.", CYAN),
    ("Disease Outbreak", "Surveillance gap demands AI prediction.", RGBColor(34, 197, 94)), # Green
    ("Rural Digital Leap", "Capture 97% offline via voice-AI now.", RGBColor(139, 92, 246)), # Purple
    ("Zoonotic Risk", "One Health integration urgently needed.", RGBColor(239, 68, 68)) # Red
]

# 2 Rows of 3
for i, (head, body, color) in enumerate(reasons):
    r = i // 3
    c = i % 3
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5 + c*4.2), Inches(2 + r*2.5), Inches(4), Inches(2.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = color
    shape.line.width = Pt(3)
    
    tf = shape.text_frame
    p = tf.add_paragraph()
    p.text = head
    set_text_style(p, 18, True, color)
    p = tf.add_paragraph()
    p.text = body
    set_text_style(p, 14, False, DARK_GREY)


# --- SLIDE 6: Hybrid Solutions Architecture ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header_strip(slide, prs.slide_width)
add_company_badge(slide)

tb = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(6), Inches(4))
tf = tb.text_frame
tf.word_wrap = True
p = tf.add_paragraph()
p.text = "Hybrid-Enabled Solutions Architecture"
set_text_style(p, 36, True, BLUE)
p = tf.add_paragraph()
p.text = "Synergizing Physical Infrastructure with Digital Intelligence to Transform Animal Care."
set_text_style(p, 24, False, DARK_GREY)

p = tf.add_paragraph()
p.text = "\nLeveraging Agentic AI, ML & Digitization to bridge the last-mile gap."
set_text_style(p, 20, True, CYAN)

img_stream = fetch_image(IMG_HYBRID)
if img_stream:
    slide.shapes.add_picture(img_stream, Inches(7), Inches(1.5), Inches(6), Inches(4.5))


# --- SLIDE 7: 10 Solutions Matrix ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header_strip(slide, prs.slide_width)
add_company_badge(slide)

create_title_box(slide, "10 Impactful Categories for Investors", x=0.5, y=0.8)

# 2 Cols, 5 Rows text table
solutions = [
    "1. Agentic AI Triage & Remote Care: Delivers autonomous 22-language triage, urgency detection, and seamless vet escalation with proactive follow-ups.",
    "2. ABDM-Compliant Digital Health Spine: Creates secure digital IDs with OCR/voice capture linking treatments, vaccines, and health history nationwide.",
    "3. AtoZ Marketplace & AI Counterfeit Shield: Authenticates products, detects counterfeit patterns, scores vendors, and ensures safe medicine/feed distribution.",
    "4. Hybrid Mobile Vans & Diagnostics: Provides AI-routed vans for X-ray, USG, labs, and home vaccinations across underserved regions.",
    "5. Clinic OS & Revenue Intelligence: Unifies queues, billing, tele-vet, labs, inventory, and real-time performance dashboards.",
    "6. Smart Supply Chain & Auto-Replenishment: Predicts demand, triggers auto-restocking, authenticates suppliers, and optimizes last-mile delivery.",
    "7. Subsidy, DBT & Welfare Intelligence: Auto-matches eligibility, files claims, issues eVouchers, and links subsidies to nutrition and disease prevention.",
    "8. Women-First Voice AI Companion: Provides voice-only guidance in 22 Indic languages for livestock health, breeding, first-aid, and nutrition.",
    "9. Predictive Wellness & Carbon Intel: Predicts diseases, models methane reduction, and qualifies livestock for carbon credit benefits.",
    "10. Community Safety & Rehoming: Uses facial recognition, geo-alerts, and verified workflows to connect lost pets and enable safe adoptions."
]

for i, sol in enumerate(solutions):
    col = i // 5 # 0 or 1
    row_in_col = i % 5
    
    x_pos = 0.5 + col * 6.5
    y_pos = 1.8 + row_in_col * 1.1
    
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x_pos), Inches(y_pos), Inches(6), Inches(1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_GREY
    shape.line.fill.background()
    
    tf = shape.text_frame
    tf.margin_left = Inches(0.1)
    tf.margin_top = Inches(0.1)
    
    parts = sol.split(": ", 1)
    title_text = parts[0]
    desc_text = parts[1]
    
    p = tf.add_paragraph()
    p.text = title_text
    set_text_style(p, 14, True, BLUE)
    
    p = tf.add_paragraph()
    p.text = desc_text
    set_text_style(p, 11, False, DARK_GREY)


# --- SLIDE 8: Vision & Mission ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header_strip(slide, prs.slide_width)
add_company_badge(slide)

img_stream = fetch_image(IMG_VISION)
if img_stream:
    slide.shapes.add_picture(img_stream, Inches(7.5), Inches(1.5), Inches(5.5), Inches(5))

tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(6.5), Inches(4.5))
tf = tb.text_frame
p = tf.add_paragraph()
p.text = "Our North Star"
set_text_style(p, 36, True, BLUE)

p = tf.add_paragraph()
p.text = "Vision"
set_text_style(p, 24, True, BLUE)
p = tf.add_paragraph()
p.text = "A trusted, AI-powered AtoZ companion for every Indian pet and livestock family, safeguarding health, livelihoods, and happiness.\n"
set_text_style(p, 18, False, DARK_GREY)

p = tf.add_paragraph()
p.text = "Mission"
set_text_style(p, 24, True, CYAN)
p = tf.add_paragraph()
p.text = "Digitize, connect, and empower care journeys with tele-vet, mobile diagnostics, and a transparent, affordable marketplace."
set_text_style(p, 18, False, DARK_GREY)


# --- SLIDE 9: Core Purpose ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header_strip(slide, prs.slide_width)
add_company_badge(slide)

create_title_box(slide, "Our Core Purpose", x=0.5, y=0.8, w=12, h=1)
slide.shapes[2].text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# 3 Columns
purposes = [
    ("Who We Are", "A purpose-driven health-tech company uniting vets, data, and marketplaces for India’s pets and livestock.", BLUE),
    ("What We Do", "Provide AI triage, remote care, and curated services so families access the right expert, product, or subsidy instantly.", CYAN),
    ("Why We Do It", "To remove distance, distrust, and complexity, turning every animal interaction into safer care, fair value, and dignity for owners.", ORANGE)
]

for i, (head, body, color) in enumerate(purposes):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5 + i*4.2), Inches(2.5), Inches(4), Inches(4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = color
    shape.line.width = Pt(3)
    
    tf = shape.text_frame
    tf.margin_top = Cm(1.0)
    
    p = tf.add_paragraph()
    p.text = head
    p.alignment = PP_ALIGN.CENTER
    set_text_style(p, 24, True, color)
    
    p = tf.add_paragraph()
    p.text = body
    p.alignment = PP_ALIGN.CENTER
    set_text_style(p, 16, False, DARK_GREY)


# --- SLIDE 10: The PetLiv Promise ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header_strip(slide, prs.slide_width)
add_company_badge(slide)

tb = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(6), Inches(4))
tf = tb.text_frame
tf.word_wrap = True
p = tf.add_paragraph()
p.text = "The PetLiv Promise"
set_text_style(p, 36, True, BLUE)

p = tf.add_paragraph()
p.text = '"AI-powered tele-vet, mobile diagnostics, and trusted marketplace give every Indian pet and livestock parent affordable, best-in-class care, products, and lifelong digital confidence."'
set_text_style(p, 24, False, DARK_GREY)

img_stream = fetch_image(IMG_PROMISE)
if img_stream:
    slide.shapes.add_picture(img_stream, Inches(7), Inches(1.5), Inches(6), Inches(4.5))


# --- SLIDE 11: 10 Ground Problem Transformations ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header_strip(slide, prs.slide_width)
add_company_badge(slide)

create_title_box(slide, "10 Ground Problem Transformations", x=0.5, y=0.5)

transformations = [
    ("Unified Records", "Zero Repeat Diagnostics", "Turn scattered animal histories into a single intelligent health identity.", BLUE),
    ("One-Stop Marketplace", "3-Hour Weekly Savings", "Replace 15 fragmented apps with one trusted commerce ecosystem.", CYAN),
    ("Counterfeit Shield", "20% Mortality Cut", "AI that protects every animal from fake medicines and toxic feed.", ORANGE),
    ("Auto-Reminders", "55% Abandonment Eliminated", "Never miss a dose again—AI ensures treatments finish, every time.", RGBColor(34, 197, 94)),
    ("Subsidy Engine", "₹10K Cr Unlocked", "Bring government benefits to every farmer with one automated tap.", BLUE),
    ("Women Voice AI", "3X Rural Retention", "Empower rural women caregivers with hands-free, language-first AI.", CYAN),
    ("Offline Digitization", "97% AI Visibility", "Digitize every animal instantly—any paper, any phone, anywhere.", ORANGE),
    ("Mobile Vet Network", "60% Emergency Saves", "AI-routed medical vans that deliver life-saving care to every village.", RGBColor(34, 197, 94)),
    ("Multilingual Content", "75% Inclusion", "Make expert livestock knowledge speak every Indian language.", BLUE),
    ("Carbon Intelligence", "Green Finance Access", "Turn healthier livestock into carbon credits and climate rewards.", CYAN)
]

for i, (p, t, d, c) in enumerate(transformations):
    col = i // 5
    row_in_col = i % 5
    
    x_pos = 0.5 + col * 6.4
    y_pos = 1.5 + row_in_col * 1.1
    
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x_pos), Inches(y_pos), Inches(6), Inches(1.0))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = c
    shape.line.width = Pt(1.5)
    
    tf = shape.text_frame
    p_title = tf.add_paragraph()
    p_title.text = f"{p} → {t}"
    set_text_style(p_title, 14, True, c)
    
    p_tagline = tf.add_paragraph()
    p_tagline.text = f'"{d}"'
    set_text_style(p_tagline, 11, False, DARK_GREY)


# --- SLIDE 12: Existing vs Future Comparison ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header_strip(slide, prs.slide_width)
add_company_badge(slide)

create_title_box(slide, "Existing System vs. Future With PetLiv AtoZ Mart AI", x=0.5, y=0.5)

rows = 10
cols = 3
table_shape = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.5))
tbl = table_shape.table

# Headers
headers = ["Category", "Existing System", "Future With PetLiv AtoZ Mart AI"]
for i, h in enumerate(headers):
    cell = tbl.cell(0, i)
    cell.text = h
    cell.fill.solid()
    cell.fill.fore_color.rgb = DARK_GREY
    p = cell.text_frame.paragraphs[0]
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.size = Pt(12)

comparison_data = [
    ("Unified Digital Records", "78% livestock & 38M pets lack digital IDs; repeat diagnostics increase costs.", "ABDM-compliant digital IDs with OCR/voice capture; 97% digitization saves ₹8K/animal/year."),
    ("One-Stop Marketplace", "Owners use 15+ fragmented apps; 3–5 hours/week wasted searching services.", "Unified marketplace with 10K+ verified vendors; instant geo-matched bookings for all needs."),
    ("Counterfeit Supply Chain", "30% fake meds/feed cause 20% mortality in ₹15K Cr industry.", "AI counterfeit detection + blockchain traceability ensuring full authenticity & safety."),
    ("Treatment Abandonment", "55% owners miss doses; ₹15K Cr annual losses due to incomplete treatments.", "ML-driven personalized reminders achieve 98.4% treatment completion."),
    ("Govt Scheme Access", "<20% access ₹10K Cr AHIDF/NLM subsidies due to paperwork complexity.", "AI eligibility-matching, auto-filing & eVouchers boost utilization to 85%."),
    ("Rural Voice Exclusion", "70% low-literacy rural women excluded by English-only apps.", "Voice-first 22-language Indic LLM increases rural retention 3X."),
    ("Offline Livestock & Pets", "97% livestock & 60% rural pets undigitized; invisible to AI systems.", "Pan-India ABDM ID coverage enabling real-time tracking for 574M animals."),
    ("Vet Access Barrier", "1:5000 vet ratio; long travel & 45-min waits cause 60% emergency deaths.", "Agentic AI triage under 4 min; tele-vet escalation + mobile vans with 95% utilization."),
    ("Language Content Gap", "90% content English-only, excluding 75% of native-language users.", "22-language knowledge hub + multilingual AI triage expanding inclusive access.")
]

for r_idx, row_data in enumerate(comparison_data):
    for c_idx, val in enumerate(row_data):
        cell = tbl.cell(r_idx+1, c_idx)
        cell.text = val
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(10)
        p.vertical_anchor = MSO_ANCHOR.MIDDLE
        p.word_wrap = True
        
        if c_idx == 0:
            p.font.bold = True
        elif c_idx == 2:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(236, 254, 255) # Light Cyan


# --- SLIDE 13: Market Dashboard ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header_strip(slide, prs.slide_width)
add_company_badge(slide)

create_title_box(slide, "Market Opportunity & Growth Landscape", x=0.5, y=0.5)

# Simplified Market Chart Placeholder
chart_placeholder = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.5), Inches(8), Inches(5.5))
chart_placeholder.fill.solid()
chart_placeholder.fill.fore_color.rgb = LIGHT_GREY
chart_placeholder.line.fill.background()

chart_tb = chart_placeholder.text_frame
chart_tb.word_wrap = True
p = chart_tb.add_paragraph()
p.text = "TAM, SAM, SOM Growth Chart Placeholder"
set_text_style(p, 20, True, DARK_GREY)
p.alignment = PP_ALIGN.CENTER

# Regional Split Placeholder
region_placeholder = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9), Inches(1.5), Inches(3.8), Inches(2.5))
region_placeholder.fill.solid()
region_placeholder.fill.fore_color.rgb = RGBColor(240, 249, 255) # Light Blue
region_placeholder.line.color.rgb = BLUE
region_placeholder.line.width = Pt(1)

p = region_placeholder.text_frame.add_paragraph()
p.text = "Regional Split & Key Drivers"
set_text_style(p, 16, True, BLUE)

p = region_placeholder.text_frame.add_paragraph()
p.text = "North (35%), South (25%), West (20%), East (15%)\nDrivers: Pet Humanization, NDLM, Carbon Credits"
set_text_style(p, 12, False, DARK_GREY)

# Growth Drivers Placeholder
drivers_placeholder = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9), Inches(4.5), Inches(3.8), Inches(2.5))
drivers_placeholder.fill.solid()
drivers_placeholder.fill.fore_color.rgb = RGBColor(255, 247, 237) # Light Orange
drivers_placeholder.line.color.rgb = ORANGE
drivers_placeholder.line.width = Pt(1)

p = drivers_placeholder.text_frame.add_paragraph()
p.text = "Growth Drivers: 20% Pet CAGR, NDLM 11.7% CAGR, AI Supply Chain."
set_text_style(p, 16, True, ORANGE)


# --- SLIDE 14: Revenue Model ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header_strip(slide, prs.slide_width)
add_company_badge(slide)

create_title_box(slide, "Revenue Model: Sustainable Growth Engines", x=0.5, y=0.5)

revenue_streams = [
    ("Marketplace Commissions", "Fees from vendors on transactions.", BLUE, "Payment & Fraud Ops"),
    ("Clinic OS SaaS", "Subscription for clinic management and analytics dashboards.", CYAN, "Cloud Infra & Dev"),
    ("Vendor Subscriptions", "Premium listings, counterfeit protection, and inventory tools.", BLUE, "Support & Marketing"),
    ("Subsidy Processing Fees", "Processing government subsidies and eVoucher distribution.", ORANGE, "Compliance & Integration"),
    ("Supply Chain Logistics", "Fees for 24-hr rural delivery and AI-driven forecasting.", CYAN, "Last-Mile Delivery"),
    ("Premium Analytics", "Subscription for multi-species health dashboards and carbon credit insights.", BLUE, "Data Modeling"),
    ("Mobile Vet Network", "Service fees from mobile diagnostic van operations (X-ray, USG, Labs).", ORANGE, "Van Maintenance & Staffing")
]

for i, (name, desc, color, cost) in enumerate(revenue_streams):
    x_pos = Inches(0.5 + (i % 2) * 6.5)
    y_pos = Inches(1.5 + (i // 2) * 1.5)
    width = Inches(6)
    height = Inches(1.4)
    
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, y_pos, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = color
    shape.line.width = Pt(2)
    
    tf = shape.text_frame
    p = tf.add_paragraph()
    p.text = name
    set_text_style(p, 16, True, color)
    
    p = tf.add_paragraph()
    p.text = desc
    set_text_style(p, 11, False, DARK_GREY)
    
    p = tf.add_paragraph()
    p.text = f"Cost Driver: {cost}"
    set_text_style(p, 10, False, RGBColor(100, 116, 139)) # Slate Gray


# --- SLIDE 15: GTM Strategy & Roadmap ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header_strip(slide, prs.slide_width)
add_company_badge(slide)

create_title_box(slide, "Overall GTM Strategy & Phased Rollout", x=0.5, y=0.5)

# GTM Model Box
model_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.5), Inches(4.5), Inches(0.5))
model_box.fill.solid()
model_box.fill.fore_color.rgb = RGBColor(239, 246, 255)
model_box.line.color.rgb = BLUE
p = model_box.text_frame.add_paragraph()
p.text = "B2B2C Model: Win Intermediaries to Onboard Families"
set_text_style(p, 14, True, BLUE)
p.alignment = PP_ALIGN.CENTER

# Segmentation (Left Column)
segmentation_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.3), Inches(4.5), Inches(5))
tf = segmentation_box.text_frame
tf.word_wrap = True

segments = [
    ("Urban Pet GTM", "Clinic-first onboarding, co-branded offers, free AI Triage."),
    ("Rural Livestock GTM", "Voice-first onboarding via FPOs/SHGs, NDLM integration."),
    ("Institutional GTM", "Position as NDLM private innovation layer, pilot subsidy auto-claims.")
]

for head, body in segments:
    p = tf.add_paragraph()
    p.text = head
    set_text_style(p, 16, True, ORANGE)
    p = tf.add_paragraph()
    p.text = body
    set_text_style(p, 12, False, DARK_GREY)
    tf.add_paragraph() # Spacer

# Roadmap (Right Column)
roadmap_box = slide.shapes.add_textbox(Inches(6), Inches(1.5), Inches(6.8), Inches(5.5))
p = roadmap_box.text_frame.add_paragraph()
p.text = "Technology Readiness & Expansion Roadmap"
set_text_style(p, 20, True, CYAN)

phases = [
    (1, "Phase 1: Pilot & Validate (TRL 7)", "Bengaluru Urban (Pet) & Birbhum (Livestock). Validate Triage, Digital ID, Clinic OS MVP.", BLUE),
    (2, "Phase 2: Deepen & Scale (TRL 8)", "Expand to 3+ cities. Launch Mobile Vans and Marketplace logistics. Deepen Govt collaboration.", CYAN),
    (3, "Phase 3: Nationwide Expansion (TRL 9)", "State-wise expansion. Full integration of Carbon Intelligence and institutional SaaS ecosystem.", ORANGE)
]

# Drawing simple roadmap dots and text
y_start = 2.5
for num, head, body, color in phases:
    # Circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6), Inches(y_start + num*1.5), Inches(0.3), Inches(0.3))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    
    # Text Box
    text_box = slide.shapes.add_textbox(Inches(6.5), Inches(y_start + num*1.5 - 0.2), Inches(6), Inches(1))
    p = text_box.text_frame.add_paragraph()
    p.text = head
    set_text_style(p, 14, True, color)
    p = text_box.text_frame.add_paragraph()
    p.text = body
    set_text_style(p, 10, False, DARK_GREY)


# --- SLIDE 16: Focused Pilot Strategy ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header_strip(slide, prs.slide_width)
add_company_badge(slide)

create_title_box(slide, "Focused Pilot Launch Strategy (Phase 1)", x=0.5, y=0.5)

# Left: Urban Pet Pilot
urban_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6), Inches(5.5))
tf = urban_box.text_frame

p = tf.add_paragraph()
p.text = "Bengaluru Urban (Pet-Focused Pilot)"
set_text_style(p, 20, True, BLUE)

p = tf.add_paragraph()
p.text = "Positioning: “Free AI vet triage and digital health IDs for all pets; smarter clinics and trusted marketplace.”"
set_text_style(p, 14, False, DARK_GREY)

p = tf.add_paragraph()
p.text = "\nKey Partners:"
set_text_style(p, 16, True, CYAN)

p = tf.add_paragraph()
p.text = "1. 10–20 Clinics (Koramangala, Whitefield, etc.)\n2. 20–30 Pet Vendors (Shops, groomers, trainers)"
set_text_style(p, 14, False, DARK_GREY)

p = tf.add_paragraph()
p.text = "\nGTM Tactics:"
set_text_style(p, 16, True, ORANGE)

p = tf.add_paragraph()
p.text = "— Deploy Clinic OS (Queue, Tele-vet) in pilot clinics.\n— In-clinic assets: standees, QR posters, Rx slips.\n— Offer free triage and digital records to all walk-in pet parents."
set_text_style(p, 14, False, DARK_GREY)

# Right: Rural Livestock Pilot
rural_box = slide.shapes.add_textbox(Inches(7), Inches(1.5), Inches(6), Inches(5.5))
tf = rural_box.text_frame

p = tf.add_paragraph()
p.text = "Birbhum, West Bengal (Livestock-Focused Pilot)"
set_text_style(p, 20, True, BLUE)

p = tf.add_paragraph()
p.text = "Positioning: “Empowering farmers with voice-AI care, instant subsidy access, and connected veterinary support.”"
set_text_style(p, 14, False, DARK_GREY)

p = tf.add_paragraph()
p.text = "\nKey Partners:"
set_text_style(p, 16, True, CYAN)

p = tf.add_paragraph()
p.text = "1. 5–10 FPOs/Co-ops & SHG networks\n2. State AH Dept for vaccination drives\n3. Local Paravets (Pashu Sakhis) as ground agents"
set_text_style(p, 14, False, DARK_GREY)

p = tf.add_paragraph()
p.text = "\nGTM Tactics:"
set_text_style(p, 16, True, ORANGE)

p = tf.add_paragraph()
p.text = "— Voice-first onboarding via Indic LLM bot (Bengali).\n— Camp activation: Village meetings & vaccination camps.\n— Incentives for Pashu Sakhis per animal registration."
set_text_style(p, 14, False, DARK_GREY)


# --- SLIDE 17: Funding Request ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header_strip(slide, prs.slide_width)
add_company_badge(slide)

create_title_box(slide, "Strategic Capital Allocation", x=0.5, y=0.5, w=12, h=1)
slide.shapes[2].text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Right: The Ask & Detail
ask_box = slide.shapes.add_textbox(Inches(7.5), Inches(1.8), Inches(5), Inches(5))
tf = ask_box.text_frame

p = tf.add_paragraph()
p.text = "Funding Request"
set_text_style(p, 24, True, BLUE)
p = tf.add_paragraph()
p.text = "₹5 Crore"
set_text_style(p, 48, True, BLUE)

p = tf.add_paragraph()
p.text = "To fuel Phase 1 Pilot & Platform Validation"
set_text_style(p, 18, False, DARK_GREY)

# Left: Visualization Chart Placeholder (For simplicity, using text description)
chart_placeholder = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(6.5), Inches(5.5))
chart_placeholder.fill.solid()
chart_placeholder.fill.fore_color.rgb = RGBColor(239, 246, 255)
chart_placeholder.line.color.rgb = BLUE
chart_placeholder.line.width = Pt(2)

p = chart_placeholder.text_frame.add_paragraph()
p.text = "Use of Funds (Visualized)"
set_text_style(p, 20, True, BLUE)

fund_details = [
    ("Sales, Marketing & Distribution", "20%", BLUE),
    ("Marketplace Onboarding & Vendor Success", "20%", CYAN),
    ("Product Development & Engineering", "20%", ORANGE),
    ("Pet & Livestock Onboarding Programs", "10%", RGBColor(34, 197, 94)),
    ("Cloud Infrastructure & Maintenance", "10%", RGBColor(139, 92, 246)),
    ("AI/ML Data & Research", "10%", RGBColor(234, 179, 8)),
    ("Operations, Compliance & Misc", "10%", RGBColor(100, 116, 139))
]

for name, pct, color in fund_details:
    p = chart_placeholder.text_frame.add_paragraph()
    run = p.add_run()
    run.text = f"{pct} {name}"
    run.font.color.rgb = color
    run.font.bold = True
    run.font.size = Pt(14)


# --- SLIDE 18: Core Assumptions ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header_strip(slide, prs.slide_width)
add_company_badge(slide)

create_title_box(slide, "Core Assumptions: High-Growth Scale", x=0.5, y=0.5)

# Left: User & Partner Growth Table
table_data_user = [
    ["Year", "Parents", "Partners"],
    ["2027", "5,000", "900"],
    ["2028", "10,000", "1,800"],
    ["2029", "20,000", "3,600"],
    ["2030", "40,000", "7,200"]
]
table_shape_user = slide.shapes.add_table(5, 3, Inches(0.5), Inches(1.5), Inches(5), Inches(2.5))
tbl_user = table_shape_user.table
for r_idx, row in enumerate(table_data_user):
    for c_idx, val in enumerate(row):
        cell = tbl_user.cell(r_idx, c_idx)
        cell.text = val
        if r_idx == 0:
            cell.fill.solid(); cell.fill.fore_color.rgb = BLUE
            p = cell.text_frame.paragraphs[0]; p.font.color.rgb = WHITE
            p.font.size = Pt(12)
        else:
            p = cell.text_frame.paragraphs[0]; p.font.size = Pt(12)

# Right: Financial Mechanics
mechanics_box = slide.shapes.add_textbox(Inches(6), Inches(1.5), Inches(6.8), Inches(5))
tf = mechanics_box.text_frame
tf.word_wrap = True

p = tf.add_paragraph()
p.text = "Revenue Mechanics"
set_text_style(p, 20, True, CYAN)

data_rows = [
    ("Txn Frequency", "10 Paid Txns / Year"),
    ("Marketplace Fee", "₹10 per Txn (Provider Paid)"),
    ("SaaS Subscription", "₹1,000 / Month / Partner")
]

for label, value in data_rows:
    p = tf.add_paragraph()
    p.text = f"- {label}: {value}"
    set_text_style(p, 14, False, DARK_GREY)

p = tf.add_paragraph()
p.text = "\nOperating Expenses"
set_text_style(p, 20, True, ORANGE)

data_rows = [
    ("2026 Build Phase (Capex)", "₹5.0 Cr"),
    ("2027 OpEx Base", "₹1.50 Cr"),
    ("OpEx Growth Rate", "10% YoY (Lean Scaling)")
]
for label, value in data_rows:
    p = tf.add_paragraph()
    p.text = f"- {label}: {value}"
    set_text_style(p, 14, False, DARK_GREY)


# --- SLIDE 19: Financial Projections ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header_strip(slide, prs.slide_width)
add_company_badge(slide)

create_title_box(slide, "Financial Projections (2027 - 2030)", x=0.5, y=0.5)

# P&L Table
table_data_pl = [
    ["Metric", "2027", "2028 (Ops Breakeven)", "2029", "2030"],
    ["Marketplace Revenue", "₹5.0 L", "₹10.0 L", "₹20.0 L", "₹40.0 L"],
    ["SaaS Revenue", "₹1.08 Cr", "₹2.16 Cr", "₹4.32 Cr", "₹8.64 Cr"],
    ["Total Revenue", "₹1.13 Cr", "₹2.26 Cr", "₹4.52 Cr", "₹9.04 Cr"],
    ["Total OpEx", "₹1.50 Cr", "₹1.65 Cr", "₹1.82 Cr", "₹2.00 Cr"],
    ["EBITDA", "-₹37.0 L", "+₹61.0 L", "+₹2.71 Cr", "+₹7.04 Cr"]
]
table_shape_pl = slide.shapes.add_table(6, 5, Inches(0.5), Inches(1.5), Inches(12.333), Inches(2.5))
tbl_pl = table_shape_pl.table

for r_idx, row_data in enumerate(table_data_pl):
    for c_idx, val in enumerate(row_data):
        cell = tbl_pl.cell(r_idx, c_idx)
        cell.text = val
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(12)
        p.alignment = PP_ALIGN.RIGHT if c_idx > 0 else PP_ALIGN.LEFT
        
        if r_idx == 0:
            cell.fill.solid(); cell.fill.fore_color.rgb = DARK_GREY
            p.font.color.rgb = WHITE
            p.font.bold = True
        elif r_idx == 5: # EBITDA row
            p.font.bold = True
            if "-" in val: p.font.color.rgb = RGBColor(239, 68, 68)
            else: p.font.color.rgb = RGBColor(34, 197, 94)
        elif r_idx == 3: # Total Revenue row
             p.font.bold = True

# Breakeven Callout
breakeven_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9), Inches(4.5), Inches(3.833), Inches(2.5))
breakeven_box.fill.solid()
breakeven_box.fill.fore_color.rgb = RGBColor(239, 246, 255) # Light Blue
breakeven_box.line.color.rgb = BLUE
breakeven_box.line.width = Pt(2)

p = breakeven_box.text_frame.add_paragraph()
p.text = "Operational Breakeven"
set_text_style(p, 20, True, BLUE)
p.alignment = PP_ALIGN.CENTER

p = breakeven_box.text_frame.add_paragraph()
p.text = "2028"
set_text_style(p, 40, True, BLUE)
p.alignment = PP_ALIGN.CENTER

p = breakeven_box.text_frame.add_paragraph()
p.text = "Achieved in second commercial year."
set_text_style(p, 12, False, DARK_GREY)
p.alignment = PP_ALIGN.CENTER

# --- SLIDE 20: Thank You ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header_strip(slide, prs.slide_width)
add_company_badge(slide)

tb = slide.shapes.add_textbox(Inches(3), Inches(3), Inches(7.333), Inches(3))
tf = tb.text_frame
p = tf.add_paragraph()
p.text = "Thank You"
p.alignment = PP_ALIGN.CENTER
set_text_style(p, 60, True, BLUE)

p = tf.add_paragraph()
p.text = "Let's Build the Future of Animal Care Together."
p.alignment = PP_ALIGN.CENTER
set_text_style(p, 24, False, DARK_GREY)

p = tf.add_paragraph()
p.text = "\ncontact@aavaayahealthcare.com | www.petliv.ai"
p.alignment = PP_ALIGN.CENTER
set_text_style(p, 18, False, ORANGE)


# Save
prs.save(FILENAME)
print(f"Presentation saved as {FILENAME}")